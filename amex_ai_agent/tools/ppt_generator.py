from __future__ import annotations

from pathlib import Path
from typing import Dict

from pptx import Presentation


def run(argument: str) -> Dict[str, str]:
    summary = argument.strip() or "No summary provided"

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Fraud Model Review"
    title_slide.placeholders[1].text = "Enterprise Agent Generated"

    sections = [
        "Model performance",
        "Key features",
        "Alerts rationalization",
        "RCA insights",
    ]
    for section in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = section
        slide.placeholders[1].text = f"{section} summary:\n{summary[:1000]}"

    output = Path("fraud_summary.pptx")
    prs.save(output)
    return {"presentation": str(output.resolve())}
